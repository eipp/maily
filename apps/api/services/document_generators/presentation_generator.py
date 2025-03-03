"""Presentation document generator."""

import os
import logging
from typing import Dict, Any, Optional, Tuple, List
from datetime import datetime

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .base_generator import BaseGenerator

logger = logging.getLogger(__name__)

class PresentationGenerator(BaseGenerator):
    """Generator for presentation documents."""

    def __init__(self, document_base_path: str, document_base_url: str):
        """Initialize the presentation generator.

        Args:
            document_base_path: Base path for document storage
            document_base_url: Base URL for document access
        """
        super().__init__(document_base_path, document_base_url)

    async def generate(
        self,
        document_id: str,
        document_data: Dict[str, Any],
        template: Optional[Dict[str, Any]] = None
    ) -> Tuple[str, str, str]:
        """Generate a presentation document.

        Args:
            document_id: Document ID
            document_data: Document data
            template: Optional template data

        Returns:
            Tuple of (file_path, file_url, preview_url)
        """
        # Generate presentation
        pptx_path = self._generate_presentation(document_id, document_data, template)
        
        # Define file paths
        file_name = f"{document_id}.pptx"
        file_path = os.path.join(self.document_base_path, file_name)
        file_url = f"{self.document_base_url}/{file_name}"
        preview_url = f"{self.document_base_url}/previews/{document_id}_preview.png"
        
        # Create preview placeholder
        self.create_preview_placeholder(document_id)
        
        # Track metrics
        self.track_metrics(file_path, "presentation", template)
        
        return file_path, file_url, preview_url

    def _generate_presentation(
        self,
        document_id: str,
        document_data: Dict[str, Any],
        template: Optional[Dict[str, Any]] = None
    ) -> str:
        """Generate a PowerPoint presentation.

        Args:
            document_id: Document ID
            document_data: Document data
            template: Optional template data

        Returns:
            Path to generated presentation file
        """
        # Create a new presentation
        prs = Presentation()
        
        # Get presentation data
        title = document_data.get("title", "Untitled Presentation")
        subtitle = document_data.get("subtitle", "")
        author = document_data.get("author", "")
        slides = document_data.get("slides", [])
        
        # Apply template if provided
        if template:
            if not title and "title" in template:
                title = template["title"]
                
            if not subtitle and "subtitle" in template:
                subtitle = template["subtitle"]
                
            if not author and "author" in template:
                author = template["author"]
                
            if not slides and "slides" in template:
                slides = template["slides"]
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
        
        # Add author and date to subtitle if provided
        if author:
            if subtitle:
                subtitle_shape.text += f"\n\n{author}"
            else:
                subtitle_shape.text = author
                
        # Add date
        today = datetime.now().strftime("%B %d, %Y")
        if subtitle_shape.text:
            subtitle_shape.text += f"\n{today}"
        else:
            subtitle_shape.text = today
        
        # Create content slides
        if not slides:
            # Create a default slide if none provided
            slides = [{
                "title": "Sample Slide",
                "content": "Add your content here"
            }]
        
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", "")
            slide_type = slide_data.get("type", "content")
            
            if slide_type == "title":
                # Title slide
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]
                
                title_shape.text = slide_title
                if slide_content:
                    subtitle_shape.text = slide_content
            
            elif slide_type == "section":
                # Section header slide
                slide_layout = prs.slide_layouts[2]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                
                title_shape.text = slide_title
            
            elif slide_type == "two_content":
                # Two content slide
                slide_layout = prs.slide_layouts[3]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                left_shape = slide.placeholders[1]
                right_shape = slide.placeholders[2]
                
                title_shape.text = slide_title
                
                # Split content for left and right
                if isinstance(slide_content, list) and len(slide_content) >= 2:
                    left_content = slide_content[0]
                    right_content = slide_content[1]
                else:
                    # Split content in half if it's a string
                    if isinstance(slide_content, str):
                        parts = slide_content.split("\n\n")
                        mid = len(parts) // 2
                        left_content = "\n\n".join(parts[:mid])
                        right_content = "\n\n".join(parts[mid:])
                    else:
                        left_content = "Left content"
                        right_content = "Right content"
                
                left_shape.text = left_content
                right_shape.text = right_content
            
            elif slide_type == "comparison":
                # Comparison slide
                slide_layout = prs.slide_layouts[3]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                left_shape = slide.placeholders[1]
                right_shape = slide.placeholders[2]
                
                title_shape.text = slide_title
                
                # Get comparison data
                comparison = slide_data.get("comparison", {})
                left_title = comparison.get("left_title", "Option A")
                right_title = comparison.get("right_title", "Option B")
                left_points = comparison.get("left_points", [])
                right_points = comparison.get("right_points", [])
                
                # Format left content
                left_text = f"{left_title}\n\n"
                for point in left_points:
                    left_text += f"• {point}\n"
                
                # Format right content
                right_text = f"{right_title}\n\n"
                for point in right_points:
                    right_text += f"• {point}\n"
                
                left_shape.text = left_text
                right_shape.text = right_text
            
            elif slide_type == "image":
                # Image slide
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                
                title_shape.text = slide_title
                
                # Add image if path provided
                image_path = slide_data.get("image_path")
                if image_path and os.path.exists(image_path):
                    left = Inches(2.5)
                    top = Inches(2)
                    width = Inches(5)
                    slide.shapes.add_picture(image_path, left, top, width=width)
                
                # Add caption if provided
                caption = slide_data.get("caption")
                if caption:
                    left = Inches(2)
                    top = Inches(5)
                    width = Inches(6)
                    height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = caption
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            else:
                # Default content slide
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                content_shape = slide.placeholders[1]
                
                title_shape.text = slide_title
                
                # Format content as bullet points if it's a list
                if isinstance(slide_content, list):
                    content_shape.text = ""
                    for i, point in enumerate(slide_content):
                        p = content_shape.text_frame.add_paragraph()
                        if i > 0:  # First paragraph already exists
                            p.text = point
                            p.level = 0
                        else:
                            content_shape.text = point
                else:
                    content_shape.text = slide_content
        
        # Save presentation
        file_name = f"{document_id}.pptx"
        file_path = os.path.join(self.document_base_path, file_name)
        
        # Ensure directory exists
        self.ensure_directories(os.path.dirname(file_path))
        
        # Save presentation
        prs.save(file_path)
        
        return file_path
